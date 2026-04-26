"""Build the pitch deck for Dawa Dost @ Bharat Bricks IIT Madras (winner).

Mapped to the rubric (Databricks 30%, Accuracy 25%, Innovation 25%, Demo 20%).
Outputs DawaDost_pitch.pptx in repo root.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ---------- design tokens ----------
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x66, 0x77)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x15, 0x65, 0xC0)
LIGHT = RGBColor(0xEC, 0xEF, 0xF4)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_chip(slide, x, y, label, *, fill=ACCENT, fg=BG, size=11):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(0.32))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Emu(50000)
    tf.margin_top = tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"
    return shp


def footer(slide, slide_no, total=6):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "Dawa Dost · दवा दोस्त · Bharat Bricks Hacks 2026 · IIT Madras",
             size=9, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.3),
             f"{slide_no}/{total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------- build ----------
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ============================================================
# SLIDE 1 — TITLE / HOOK
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
# big colored sidebar
add_rect(s, 0, 0, Inches(0.4), H, NAVY)

add_text(s, Inches(0.9), Inches(0.6), Inches(11), Inches(0.45),
         "BHARAT BRICKS HACKS 2026 · IIT MADRAS",
         size=14, bold=True, color=NAVY)

add_text(s, Inches(0.9), Inches(1.2), Inches(11), Inches(2.0),
         "Dawa Dost", size=80, bold=True, color=INK)
add_text(s, Inches(0.9), Inches(2.6), Inches(11), Inches(0.7),
         "दवा दोस्त · the medicine friend", size=28, color=ACCENT)

add_text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
         "Translator apps tell a rural patient WHAT their prescription says.",
         size=22, color=MUTED)
add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
         "Dawa Dost tells them whether to TRUST it.",
         size=24, bold=True, color=INK)

# stats strip
strip_y = Inches(5.4)
strip_h = Inches(1.1)
add_rect(s, Inches(0.9), strip_y, Inches(11.5), strip_h, LIGHT)
def stat(x, num, label, color=NAVY):
    add_text(s, x, strip_y + Inches(0.10), Inches(3.7), Inches(0.55), num, size=28, bold=True, color=color)
    add_text(s, x, strip_y + Inches(0.62), Inches(3.7), Inches(0.4), label, size=11, color=MUTED)
stat(Inches(1.1), "280M", "low-literacy Indians", NAVY)
stat(Inches(5.0), "230+", "NSQ batches recalled / yr (CDSCO)", RED)
stat(Inches(8.9), "60–90%", "savings via Jan Aushadhi generics", GREEN)

footer(s, 1)


# ============================================================
# SLIDE 2 — THE PROBLEM (Innovation 25%)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, 0, W, Inches(0.7), NAVY)
add_text(s, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Why this is a real Indian problem, not a generic one",
         size=20, bold=True, color=BG)

# Two strangers narrative
add_text(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
         "Between an Indian patient and safe medicine sit two strangers:",
         size=20, color=INK)

card_w = Inches(5.9); card_h = Inches(2.2)
def card(x, y, title, body, color):
    add_rect(s, x, y, card_w, card_h, LIGHT)
    add_rect(s, x, y, card_w, Inches(0.45), color)
    add_text(s, x + Inches(0.3), y + Inches(0.05), card_w - Inches(0.6), Inches(0.4), title, size=15, bold=True, color=BG)
    add_text(s, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.6), card_h - Inches(0.7), body, size=14, color=INK)

card(Inches(0.6), Inches(1.85), "1 · The doctor's handwriting",
     "Half a billion prescriptions a year are illegible to the patient who pays for them. Translator apps render the words. They don't verify the substance.",
     NAVY)
card(Inches(6.85), Inches(1.85), "2 · The pharmacist's shelf",
     "CDSCO publishes banned drugs and Not-of-Standard-Quality batch recalls every month. The patient never sees that list. Substandard medicine looks identical to safe medicine.",
     RED)

# Insight strip
add_rect(s, Inches(0.6), Inches(4.4), Inches(12.15), Inches(1.4), NAVY)
add_text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.5),
         "The insight",
         size=14, bold=True, color=ORANGE)
add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.85),
         "The trust data is already public — sitting on a Government website nobody visits. Dawa Dost is a lakehouse problem dressed up as a healthcare problem.",
         size=18, bold=True, color=BG)

# Why us strip
add_text(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.4),
         "Bilingual by default · designed for ₹6,000 phones · works offline-first · regulatory-grade audit on every call",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

footer(s, 2)


# ============================================================
# SLIDE 3 — STACK ARCHITECTURE (mirror of README diagram)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, 0, W, Inches(0.7), NAVY)
add_text(s, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Stack architecture",
         size=20, bold=True, color=BG)
add_text(s, Inches(0.5), Inches(0.78), Inches(12), Inches(0.3),
         "Patient phone → Databricks Apps → Model Serving + Vector Search + Unity Catalog · last-mile voice via Sarvam + Bolna",
         size=11, color=MUTED)

# ---- helper for arrow lines ----
from pptx.enum.shapes import MSO_CONNECTOR
def arrow(x1, y1, x2, y2, color=NAVY, width=1.5):
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrow = True  # may be a no-op on some pptx versions

def labeled_box(x, y, w, h, title, body, fill=NAVY, fg=BG, body_color=ORANGE):
    add_rect(s, x, y, w, h, fill)
    add_text(s, x + Inches(0.18), y + Inches(0.10), w - Inches(0.36), Inches(0.4),
             title, size=13, bold=True, color=fg)
    add_text(s, x + Inches(0.18), y + Inches(0.55), w - Inches(0.36), h - Inches(0.65),
             body, size=10, color=body_color)

# Top row — patient ↔ apps
labeled_box(Inches(0.5),  Inches(1.3), Inches(2.6), Inches(1.1),
            "📱 Patient phone",
            "PWA · vanilla JS\nHindi · English · Malayalam",
            fill=LIGHT, fg=INK, body_color=MUTED)

labeled_box(Inches(4.7),  Inches(1.3), Inches(4.0), Inches(1.1),
            "Databricks Apps",
            "FastAPI · M2M OAuth\napp/main.py + routers",
            fill=ACCENT, fg=BG, body_color=BG)

labeled_box(Inches(10.3), Inches(1.3), Inches(2.6), Inches(1.1),
            "📞 Bolna · 🔊 Sarvam",
            "Indic voice agent\nBulbul v2 TTS · Saarika ASR",
            fill=LIGHT, fg=INK, body_color=MUTED)

# Arrows top row
arrow(Inches(3.1),  Inches(1.85), Inches(4.7),  Inches(1.85))   # phone → apps (request)
arrow(Inches(4.7),  Inches(2.05), Inches(3.1),  Inches(2.05))   # apps → phone (response)
arrow(Inches(8.7),  Inches(1.85), Inches(10.3), Inches(1.85))   # apps → bolna/sarvam
arrow(Inches(10.3), Inches(2.05), Inches(8.7),  Inches(2.05))   # response

# Tiny labels along arrows
add_text(s, Inches(3.15), Inches(1.55), Inches(1.5), Inches(0.25), "prescription jpg", size=9, color=MUTED)
add_text(s, Inches(3.15), Inches(2.10), Inches(1.5), Inches(0.25), "bilingual audio + JSON", size=9, color=MUTED)
add_text(s, Inches(8.75), Inches(1.55), Inches(1.5), Inches(0.25), "user_data", size=9, color=MUTED)
add_text(s, Inches(8.75), Inches(2.10), Inches(1.5), Inches(0.25), "call · transcript", size=9, color=MUTED)

# Middle band — three big Databricks rails
band_y = Inches(2.9)
labeled_box(Inches(0.5),  band_y, Inches(4.0), Inches(2.2),
            "Mosaic AI Model Serving",
            ("• llama-4-maverick — vision OCR\n"
             "• llama-3-3-70b-instruct — reasoning\n"
             "• gte-large-en — embeddings"),
            fill=NAVY)

labeled_box(Inches(4.7),  band_y, Inches(4.0), Inches(2.2),
            "Mosaic AI Vector Search",
            ("• endpoint  hack_cdsco_endpoint\n"
             "• Delta-sync  cdsco_approved_idx\n"
             "• CDC enabled · grounded RAG"),
            fill=NAVY)

labeled_box(Inches(8.9),  band_y, Inches(4.0), Inches(2.2),
            "Unity Catalog · Delta Lake",
            ("• cdsco_approved / banned / nsq_alerts / fdc\n"
             "• nlem_2022 (650) · pmbjp_catalog (2,438)\n"
             "• drug_aliases (43) · drug_food (35)\n"
             "• patient_sessions · drug_timetable\n"
             "• side_effect_log · sos_events\n"
             "• inference_log  ← MLflow-style audit"),
            fill=NAVY)

# Connect Apps → middle band
arrow(Inches(6.7), Inches(2.4), Inches(2.5),  band_y)            # apps → model serving
arrow(Inches(6.7), Inches(2.4), Inches(6.7),  band_y)            # apps → vector search
arrow(Inches(6.7), Inches(2.4), Inches(10.9), band_y)            # apps → unity catalog

# Vector Search ↔ UC retrieval
arrow(Inches(8.7), band_y + Inches(1.0), Inches(8.9), band_y + Inches(1.0))
add_text(s, Inches(8.75), band_y + Inches(0.65), Inches(0.4), Inches(0.25), "↔ retrieval", size=8, color=ORANGE)

# Bottom band — deploy + audit
deploy_y = Inches(5.4)
labeled_box(Inches(0.5),  deploy_y, Inches(6.1), Inches(1.3),
            "Deploy · Databricks Asset Bundles",
            "databricks bundle deploy   →   ingest job + Vector Search + App + secrets in one shot",
            fill=ORANGE, fg=BG, body_color=BG)

labeled_box(Inches(6.8),  deploy_y, Inches(6.1), Inches(1.3),
            "Audit · inference_log Delta table",
            "every LLM call: prompt · model · latency · tokens · citation set · session_id",
            fill=GREEN, fg=BG, body_color=BG)

footer(s, 3)


# ============================================================
# SLIDE 4 — DATABRICKS DEEP DIVE (30% — biggest weight)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, 0, W, Inches(0.7), NAVY)
add_text(s, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "How Databricks does the heavy lifting",
         size=20, bold=True, color=BG)

# Top row: client + apps
add_rect(s, Inches(0.6), Inches(1.0), Inches(2.4), Inches(0.9), LIGHT)
add_text(s, Inches(0.6), Inches(1.05), Inches(2.4), Inches(0.4), "📱 Patient phone", size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.45), Inches(2.4), Inches(0.4), "PWA · vanilla JS · Hindi/EN/ML", size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(3.4), Inches(1.0), Inches(3.0), Inches(0.9), ACCENT)
add_text(s, Inches(3.4), Inches(1.05), Inches(3.0), Inches(0.4), "Databricks Apps", size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.4), Inches(1.45), Inches(3.0), Inches(0.4), "FastAPI · M2M OAuth · Secrets", size=10, color=BG, align=PP_ALIGN.CENTER)

add_rect(s, Inches(6.8), Inches(1.0), Inches(2.6), Inches(0.9), ACCENT)
add_text(s, Inches(6.8), Inches(1.05), Inches(2.6), Inches(0.4), "Asset Bundles", size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(1.45), Inches(2.6), Inches(0.4), "one-command deploy", size=10, color=BG, align=PP_ALIGN.CENTER)

add_rect(s, Inches(9.8), Inches(1.0), Inches(3.0), Inches(0.9), LIGHT)
add_text(s, Inches(9.8), Inches(1.05), Inches(3.0), Inches(0.4), "🔊 Sarvam · 📞 Bolna", size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.8), Inches(1.45), Inches(3.0), Inches(0.4), "TTS / ASR · Indic voice agent", size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Middle: model serving (3 endpoints)
add_text(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.4),
         "Mosaic AI Model Serving — three endpoints, three jobs:", size=13, bold=True, color=NAVY)

ms_y = Inches(2.65)
def ms_card(x, model, role):
    add_rect(s, x, ms_y, Inches(4.05), Inches(1.0), NAVY)
    add_text(s, x + Inches(0.15), ms_y + Inches(0.1), Inches(3.8), Inches(0.4), model, size=12, bold=True, color=BG)
    add_text(s, x + Inches(0.15), ms_y + Inches(0.55), Inches(3.8), Inches(0.4), role, size=11, color=ORANGE)
ms_card(Inches(0.6),  "Llama-4-Maverick (vision)", "prescription OCR")
ms_card(Inches(4.75), "Llama-3-3-70B-Instruct",   "reasoning · conflict checks · explainer")
ms_card(Inches(8.9),  "GTE-Large-EN",              "embeddings for Vector Search")

# Bottom: data plane
add_text(s, Inches(0.6), Inches(3.95), Inches(12), Inches(0.4),
         "Unity Catalog · 13 Delta tables · 1 Vector Search index · 1 audit log:", size=13, bold=True, color=NAVY)

dl_y = Inches(4.45)
def dl_card(x, w, title, body, fill):
    add_rect(s, x, dl_y, w, Inches(2.1), fill)
    add_text(s, x + Inches(0.2), dl_y + Inches(0.1), w - Inches(0.4), Inches(0.4), title, size=12, bold=True, color=BG)
    add_text(s, x + Inches(0.2), dl_y + Inches(0.55), w - Inches(0.4), Inches(1.5), body, size=10, color=BG)

dl_card(Inches(0.6), Inches(3.5),
        "Reference (7 tables)",
        "CDSCO approved (51) · banned (21) · NSQ alerts (20) · FDC (20) · NLEM 2022 (650) · PMBJP (2,438) · drug-aliases (43)",
        NAVY)
dl_card(Inches(4.3), Inches(3.5),
        "Vector Search",
        "hack_cdsco_endpoint · Delta-sync index over CDSCO + NSQ · CDC enabled · embedded with gte-large-en · grounded RAG",
        ACCENT)
dl_card(Inches(8.0), Inches(4.75),
        "Operational + Audit (6 tables)",
        "patient_sessions · drug_timetable · side_effect_log · sos_events · bolna_call_outcomes · inference_log (every LLM call: prompt, model, latency, tokens, citations)",
        RED)

footer(s, 4)


# ============================================================
# SLIDE 5 — ACCURACY & SAFETY (25%)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, 0, W, Inches(0.7), NAVY)
add_text(s, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Accuracy & safety — verifiable, not vibes",
         size=20, bold=True, color=BG)

# Left: guardrails table
add_text(s, Inches(0.6), Inches(1.0), Inches(6.5), Inches(0.4),
         "7 hard guards (app/lib/guards.py)", size=15, bold=True, color=NAVY)

guards = [
    ("Dose invention", "Refused — doses only ever from prescription OCR"),
    ("Banned drug in Rx", "RAG result overridden, banned-drug alert wins"),
    ("Low-confidence RAG", "Threshold gate · says \"I don't know\" instead of guessing"),
    ("OCR brand typo", "drug_aliases exact → rapidfuzz ≥85 → LLM fallback"),
    ("Prompt injection", "Input sanitiser + output schema + banned phrases"),
    ("PII leakage", "Phone + name hashed before any logging"),
    ("Regulatory drift", "Every CDSCO row carries source_url + last_verified_at"),
]
gy = Inches(1.55)
for label, body in guards:
    add_text(s, Inches(0.6), gy, Inches(2.4), Inches(0.4), "✓ " + label, size=11, bold=True, color=GREEN)
    add_text(s, Inches(3.0), gy, Inches(4.1), Inches(0.4), body, size=10, color=INK)
    gy += Inches(0.45)

# Right: eval scorecard
right_x = Inches(7.5)
add_text(s, right_x, Inches(1.0), Inches(5.5), Inches(0.4),
         "Eval scorecard", size=15, bold=True, color=NAVY)

# big number
add_rect(s, right_x, Inches(1.55), Inches(5.5), Inches(1.4), NAVY)
add_text(s, right_x, Inches(1.65), Inches(5.5), Inches(0.6), "29 / 30", size=40, bold=True, color=BG, align=PP_ALIGN.CENTER)
add_text(s, right_x, Inches(2.4), Inches(5.5), Inches(0.5), "UI eval (Playwright, 11 screens × 3 langs)", size=11, color=ORANGE, align=PP_ALIGN.CENTER)

# eval bullets
add_text(s, right_x, Inches(3.15), Inches(5.5), Inches(0.4), "16 case-level evals (eval/cases.json)", size=12, bold=True, color=NAVY)
case_lines = [
    "  Banned-FDC detection: paracetamol+flupirtine ✓",
    "  NSQ batch flag: cefixime CXM2509A ✓",
    "  Drug-food: warfarin + green leafy ✓",
    "  Hindi explanation grounded to citation ✓",
    "  Refusal on out-of-scope question ✓",
]
for i, line in enumerate(case_lines):
    add_text(s, right_x, Inches(3.55) + Inches(0.32) * i, Inches(5.5), Inches(0.32), line, size=11, color=INK)

# audit
add_rect(s, right_x, Inches(5.4), Inches(5.5), Inches(1.4), LIGHT)
add_text(s, right_x + Inches(0.2), Inches(5.5), Inches(5.1), Inches(0.4), "MLflow-style audit log", size=12, bold=True, color=NAVY)
add_text(s, right_x + Inches(0.2), Inches(5.85), Inches(5.1), Inches(0.9),
         "Every LLM call (OCR, reasoning, explainer) appended to inference_log Delta table:\nprompt · model · latency · tokens · citation set · session_id",
         size=10, color=INK)

footer(s, 5)


# ============================================================
# SLIDE 6 — DEMO MOMENT + ASKS
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, 0, W, Inches(0.7), NAVY)
add_text(s, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "What you'll see in the demo",
         size=20, bold=True, color=BG)

# Demo timeline
flow = [
    ("①", "Scan", "OCR a real prescription via Llama-4 vision"),
    ("②", "Trust", "RAG hits CDSCO + NSQ tables; one drug flagged 🚫"),
    ("③", "Save", "₹85 → ₹12 Jan Aushadhi alternative surfaces"),
    ("④", "Speak", "Hindi voice explains the regimen via Sarvam"),
    ("⑤", "Ring", "20s countdown → Bolna calls patient in Hindi (LIVE)"),
    ("⑥", "SOS", "One red button → SMS + live call to Dr. Saab"),
]
fy = Inches(1.0)
for num, title, body in flow:
    add_text(s, Inches(0.6), fy, Inches(0.5), Inches(0.5), num, size=22, bold=True, color=ACCENT)
    add_text(s, Inches(1.2), fy + Inches(0.05), Inches(2), Inches(0.4), title, size=14, bold=True, color=INK)
    add_text(s, Inches(3.2), fy + Inches(0.08), Inches(9.5), Inches(0.4), body, size=13, color=MUTED)
    fy += Inches(0.55)

# Repo + URL strip
add_rect(s, Inches(0.6), Inches(4.5), Inches(12.15), Inches(1.5), LIGHT)
add_text(s, Inches(0.85), Inches(4.6), Inches(11.5), Inches(0.45), "Reproducible — clone + deploy in 15 minutes", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(5.05), Inches(11.5), Inches(0.4),
         "github.com/Ruchir-r/Bharatbricks_hackathon  →  README.md  →  DEPLOY.md",
         size=12, color=ACCENT, font="Consolas")
add_text(s, Inches(0.85), Inches(5.45), Inches(11.5), Inches(0.4),
         "databricks bundle deploy   ·   databricks bundle run ingest_cdsco   ·   open the app URL",
         size=11, color=INK, font="Consolas")

# Final line
add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
         "It's not a translator. It's a trust layer.",
         size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer(s, 6)


# ---------- save ----------
out = "/Users/ruchir/Desktop/claude/DawaDost_pitch.pptx"
prs.save(out)
print(f"OK · {out}")
